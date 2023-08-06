import re
from dataclasses import dataclass

import chevron
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import PicturePlaceholder
from pptx.slide import Slide, SlideLayout


def find_shape_by_name(slide: Slide, shape_name: str):
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape


def slide_layout_by_name(deck: Presentation, slide_master_name: str) -> SlideLayout:
    return next(filter(lambda l: l.name == slide_master_name, deck.slide_layouts))


def format_key_value_shape(shape):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            key_value = key_value_pattern.search(r.text)
            if key_value:
                value_run = p.add_run()

                r.text = key_value.group("key")
                r.font.bold = True
                value_run.text = key_value.group("value")


key_value_pattern = re.compile(r"(?P<key>.*:)(?P<value>.*)", re.IGNORECASE)


def format_paragraph(paragraph):
    def strip_leading_dashes(p):
        if p.runs:
            p.runs[0].text = p.runs[0].text.lstrip("- ")

    def count_leading_dashes(text: str) -> int:
        if text.startswith("-"):
            return 1 + count_leading_dashes(text[1:])
        return 0

    level = count_leading_dashes(paragraph.text)
    if level:
        paragraph.level = level
        strip_leading_dashes(paragraph)


def replace_shape_with_picture(slide: Slide, shape: BaseShape, img):
    if isinstance(shape, Picture):
        replace_with_picture(slide, shape, img)
    elif isinstance(shape, PicturePlaceholder):
        replace_placeholder_with_image(slide, shape, img)
    else:
        raise ValueError(f"""Cannot replace {type(shape).__name__} with picture""")


def update_and_format_shape_text(shape, text):
    # TODO support a subset of MarkDown format
    # TODO support of "<b>Name:</b> value" formatting
    shape.text = text
    for p in shape.text_frame.paragraphs:
        format_paragraph(p)


def is_picture_variable(variable: str) -> bool:
    return ("picture" in variable) or ("image" in variable)


def replace_with_picture(slide: Slide, picture: Picture, img):
    new_shape = slide.shapes.add_picture(
        img,
        picture.left,
        picture.top,
        picture.width,
        picture.height,
    )
    new_shape.auto_shape_type = picture.auto_shape_type
    old_pic = picture._element
    new_pic = new_shape._element
    old_pic.addnext(new_pic)
    old_pic.getparent().remove(old_pic)


def replace_placeholder_with_image(slide: Slide, placeholder: PicturePlaceholder, img):
    pic = slide.shapes.add_picture(img, placeholder.left, placeholder.top)

    # calculate max width/height for target size
    ratio = min(
        placeholder.width / float(pic.width), placeholder.height / float(pic.height)
    )

    pic.height = int(pic.height * ratio)
    pic.width = int(pic.width * ratio)

    pic.left = int(placeholder.left + ((placeholder.width - pic.width) / 2))
    pic.top = int(placeholder.top + ((placeholder.height - pic.height) / 2))

    p = placeholder.element
    p.getparent().remove(p)


def expand_variables_for_presentation(
    presentation: Presentation, variables: dict
):
    for slide in presentation.slides:
        expand_variables_on_slide(slide, variables)

    # expand core properties
    presentation.core_properties.author = expand_variables(
        presentation.core_properties.author, variables
    )
    presentation.core_properties.title = expand_variables(
        presentation.core_properties.title, variables
    )
    presentation.core_properties.subject = expand_variables(
        presentation.core_properties.subject, variables
    )


def variable_name(name: str) -> str:
    # we use mustache-line names for templating
    return f"{{{{{name}}}}}"


def expand_variables(text, variables):
    return chevron.render(text, variables)


def expand_variables_on_shape(shape, variables: dict):
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.runs and paragraph.text:
            for run in paragraph.runs:
                run.text = expand_variables(run.text, variables)

        format_paragraph(paragraph)


@dataclass
class FontSettings(object):
    font_family: str
    font_size: int
    line_height: float


default_font_settings = FontSettings("Arial", 15, 1.1)


def expand_variables_on_slide(
    slide: Slide,
    variables: dict,
    key_value_variables: list,
    autofit_variables: list,
    font_settings: FontSettings,
):
    # substitute shape values first
    for variable, value in variables.items():
        shape = find_shape_by_name(slide, variable_name(variable))
        if shape:
            if is_picture_variable(variable):
                replace_shape_with_picture(slide, shape, value)
            else:
                update_and_format_shape_text(shape, value)
                if variable in key_value_variables:
                    format_key_value_shape(shape)

                if variable in autofit_variables:
                    shape.text_frame.fit_text(
                        font_settings.font_family,
                        font_settings.font_size,
                        font_settings.line_height,
                    )

    # then expand content
    # expanding content at this point allows us to use templating in values
    for shape in slide.shapes:
        if shape.has_text_frame:
            expand_variables_on_shape(shape, variables)
        elif shape.has_table:
            for r in shape.table.rows:
                for c in r.cells:
                    expand_variables_on_shape(c, variables)


def add_slide(
    deck: Presentation,
    slide_master_name: str,
    variables: dict,
    autofit_text_variable_names: list = [],
    key_value_variable_names: list = [],
    font_settings: FontSettings = default_font_settings,
) -> Slide:
    master_layout = slide_layout_by_name(deck, slide_master_name)
    result = deck.slides.add_slide(master_layout)

    expand_variables_on_slide(
        result,
        variables,
        key_value_variable_names,
        autofit_text_variable_names,
        font_settings,
    )

    return result
