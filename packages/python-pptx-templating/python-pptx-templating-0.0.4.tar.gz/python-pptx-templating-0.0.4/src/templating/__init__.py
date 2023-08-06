# FIXME this file contains fixes to python-pptx to work properly
# in the scenarios of templating using master slides.
# Corresponding pull requests to python-pptx are in progress

from pptx.shapes.picture import Picture
from pptx.shapes.shapetree import _BaseShapes
from pptx.text.layout import TextFitter, _rendered_size


class _CloneableBaseShapes(_BaseShapes):
    def clone_placeholder(self, placeholder):
        """Add a new placeholder shape based on *placeholder*."""
        sp = placeholder.element
        ph_type, orient, sz, idx = (sp.ph_type, sp.ph_orient, sp.ph_sz, sp.ph_idx)
        id_ = self._next_shape_id
        name = self._next_name(placeholder.name, id_)
        self._spTree.add_placeholder(id_, name, ph_type, orient, sz, idx)

    def _next_name(self, basename, shape_id):
        """
        Next unique placeholder name for placeholder of name *basename*,
        with id number *shape_id*. Usually will be
        placeholder original name suffixed with id-1, e.g.
        _next_name("placeholder", 4) ==>
        'placeholder 3'. The number is incremented as necessary to make
        the name unique within the collection.
        """

        # start with basename
        name = basename

        # increment numpart as necessary to make name unique
        numpart = shape_id - 1

        names = self._spTree.xpath("//p:cNvPr/@name")
        while name in names:
            name = "%s %d" % (basename, numpart)
            numpart += 1

        return name


def _fix_placeholder_cloning():
    _BaseShapes._next_name = _CloneableBaseShapes._next_name
    _BaseShapes.clone_placeholder = _CloneableBaseShapes.clone_placeholder


class _FixedTextFitter(TextFitter):
    def _wrap_lines(self, line_source, point_size):
        """
        Return a sequence of str values representing the text in
        *line_source* wrapped within this fitter when rendered at
        *point_size*.
        """
        break_lines = self._break_line(line_source, point_size)
        if not break_lines:
            # it does not fit, no options...
            return None

        text, remainder = break_lines
        lines = [text]
        if remainder:
            wrapped_lines = self._wrap_lines(remainder, point_size)
            if not wrapped_lines:
                return None
            lines.extend(wrapped_lines)
        return lines

    @property
    def _fits_inside_predicate(self):
        """Return  function taking an integer point size argument.

        The function returns |True| if the text in this fitter can be wrapped to fit
        entirely within its extents when rendered at that point size.
        """

        def predicate(point_size):
            """Return |True| when text in `line_source` can be wrapped to fit.

            Fit means text can be broken into lines that fit entirely within `extents`
            when rendered at `point_size` using the font defined in `font_file`.
            """
            text_lines = self._wrap_lines(self._line_source, point_size)
            if not text_lines:
                # text does not fit
                return False

            # take into account line height, as it is not just font height that matters
            line_height = 1.3
            cy = _rendered_size("Ty", point_size, self._font_file)[1]
            return (line_height * cy * len(text_lines)) <= self._height

        return predicate


def _fix_text_fitting():
    TextFitter._wrap_lines = _FixedTextFitter._wrap_lines
    TextFitter._fits_inside_predicate = _FixedTextFitter._fits_inside_predicate


@property
def get_picture_description(self):
    attrib = self._element._nvXxPr.cNvPr.attrib
    if attrib.has_key("descr"):
        return attrib.get("descr")
    return None


def _add_picture_description_property():
    setattr(Picture, "description", get_picture_description)


_add_picture_description_property()
_fix_placeholder_cloning()
_fix_text_fitting()
